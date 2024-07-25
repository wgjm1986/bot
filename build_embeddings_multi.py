import sys

import docx
import pptx
import re
import csv
import os

from glob import glob
from multiprocessing import Pool
from datetime import datetime
from pypdf import PdfReader
import openai

from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_openai import OpenAIEmbeddings
from langchain_community.vectorstores import FAISS

import sqlite3
import numpy as np

from concurrent.futures import ThreadPoolExecutor

client = openai.OpenAI()

extensions = ["pdf","tex","txt","docx","pptx"]
# filenames = ['../syllabus.pdf'] \
#   + [filename for filename in glob(f"../Module 2/**/*.pdf",recursive=True)]
filenames = ['../syllabus.pdf'] \
    + [filename for ext in extensions for filename in glob(f"../Module */**/*.{ext}",recursive=True)]
#   + [filename for ext in extensions for filename in glob(f"../Discussion/**/*.{ext}",recursive=True)] \
#   + [filename for ext in extensions for filename in glob(f"../Transcripts/**/*.{ext}",recursive=True)] \
#   + [filename for ext in extensions for filename in glob(f"../Corporate finance slides/**/*.{ext}",recursive=True)] \
#   + [filename for ext in extensions for filename in glob(f"../Textbook/**/*.{ext}",recursive=True)]



print("Num files: " + str(len(filenames)))

def get_slide_text(slide):
    slide_text_chunks = []
    for shape in slide.shapes:
        if hasattr(shape,"text"):
            slide_text_chunks.append(shape.text+' ')
    return '\n'.join(slide_text_chunks)

def get_document_text(file_path):
    extension = file_path[-3:]
    if extension == "txt" or extension == "tex":
        with open(file_path, 'r') as file:
            document_text = file.read()
    elif extension == "pdf":
        document_text = ""
        pdf = PdfReader(file_path)
        for page in pdf.pages:
            try:
                document_text += page.extract_text() + "\n\n"
                # document_text += page.extract_text(extraction_mode="layout") + "\n\n"
            except:
                print("PDF read failure: " + file_path)
    elif extension == "docx":
        doc = docx.Document(file_path)
        document_text = '\n\n'.join([p.text for p in doc.paragraphs])
    elif extension == "pptx":
        pres = pptx.Presentation(file_path)
        document_text = '\n\n'.join([get_slide_text(slide) for slide in pres.slides])
    else:
        print("Unsupported file type: " + file_path)
        return
    if document_text and document_text != "":
        # print("success: " + file_path)
        if re.match("Exam|exam|Midterm|midterm",file_path):
            source_type = "Tests, midterms, and exams from past years"
        elif re.match("Module",file_path):
            source_type = "Teaching materials"
        elif re.match("Textbook",file_path):
            source_type = "Textbook"
        elif re.match("Transcripts",file_path):
            source_type = "Transcripts of class sessions"
        elif re.match("Discussion",file_path):
            source_type = "Media articles"
        else:
            source_type = "Other"
        return document_text
    else:
        print("failure: " + file_path)
        return

db_temp_path = 'my_db_tmp.db'
if os.path.exists(db_temp_path):
        os.remove(db_temp_path)
        print(f"Deleted existing database file: {db_temp_path}")
conn = sqlite3.connect(db_temp_path)
cursor = conn.cursor()
cursor.execute('''
    CREATE TABLE IF NOT EXISTS documents (
        doc_id INTEGER PRIMARY KEY,
        file_path TEXT,
        description TEXT
    )
''')
conn.commit()
cursor.execute('''
    CREATE TABLE IF NOT EXISTS chunks (
        id INTEGER PRIMARY KEY,
        doc_id INTEGER,
        chunk_text TEXT,
        embedding BLOB,
        FOREIGN KEY(doc_id) REFERENCES documents(doc_id)
    )
''')
conn.commit()
conn.close()

text_splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=50, add_start_index=True)

print(f"Importing documents: {datetime.now():%H:%M:%S}")

def process_file(filename):
    if os.path.getsize(filename) > 2e6: return
    document_text = get_document_text(filename)
    if not document_text: return
    if type(document_text) is not str: return
    chunks = text_splitter.split_text( document_text )
    chunks = [chunk for chunk in chunks if re.search('[A-Za-z]',chunk)]
    for chunk in chunks: chunk = re.sub('\t',' ',chunk)
    for chunk in chunks: chunk = re.sub(' +',' ',chunk)
    if not chunks: return
    # Get description
    description_prompt = "Please reply with a short description for the document below (30 words or fewer). Your description does not need to be a complete sentence. It should consist only of ASCII characters with no tabs, newlines, or form feeds. Be sure to mention any authors, and the year of publication, is you can find them. If the document is an exam, specify the semester, and which exam it was (first midterm, second midterm, final exam, etc.). Then at the end of the same line, list 5 or fewer keywords for the document's content."
    client = openai.OpenAI()
    description_messages = [{"role":"system","content":description_prompt} , {"role":"user","content":document_text}]
    description = client.chat.completions.create(
        model="gpt-4o-mini",
        messages=description_messages,
        max_tokens=200,
        stream=False
    ).choices[0].message.content
    # print(filename)
    # print(description)
    # Add this file to the table of files
    conn = sqlite3.connect(db_temp_path)
    cursor = conn.cursor()
    cursor.execute('''
        INSERT INTO documents (file_path, description)
        VALUES (?, ?)
    ''', (filename, description))
    doc_id = cursor.lastrowid
    for chunk in chunks:
        # Add each chunk and embeddings into the table of chunks
        embedding = openai.embeddings.create(model="text-embedding-ada-002",input=chunk).data[0].embedding
        embedding_bytes = np.array( embedding ).tobytes()
        cursor.execute('''
            INSERT INTO chunks (doc_id, chunk_text,embedding) VALUES (?, ?, ?)
        ''', (doc_id, chunk, embedding_bytes) )
    conn.commit()
    conn.close()
    print("Complete: " + filename)

with ThreadPoolExecutor(max_workers=7) as executor:
    executor.map(process_file, filenames)

print(f"Finished importing documents: {datetime.now():%H:%M:%S}")

## Now that the database construction has ended successfully, overwrite the existing one (if present)
os.rename(db_temp_path,'my_db.db')
